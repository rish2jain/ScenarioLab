"""Document ingestion and processing for seed materials."""

import logging
import uuid

from pydantic import BaseModel, Field

from app.graph.neo4j_client import get_application_neo4j_client

logger = logging.getLogger(__name__)

# In-memory cache for seed materials (backed by SQLite via SeedRepository)
_seed_store: dict[str, "SeedMaterial"] = {}

# ``delete_seed`` does not take the per-seed extraction lock; a background
# extraction task can be between ``await``s in the create_node loop. Tombstone
# only after we know the seed exists so random IDs do not leak into this set.
_tombstoned_seed_graph_ids: set[str] = set()


def tombstone_seed_graph(seed_id: str) -> None:
    """Mark a seed as deleted for in-flight graph extraction (Neo4j writes)."""
    _tombstoned_seed_graph_ids.add(seed_id)


def is_seed_graph_tombstoned(seed_id: str) -> bool:
    return seed_id in _tombstoned_seed_graph_ids


def reset_seed_graph_tombstones_for_tests() -> None:
    """Clear tombstone registry (pytest isolation)."""
    _tombstoned_seed_graph_ids.clear()


def clear_seed_graph_tombstone(seed_id: str) -> None:
    """Remove a seed from the graph tombstone set (e.g. after a failed delete retry)."""
    _tombstoned_seed_graph_ids.discard(seed_id)


class SeedGraphCleanupError(RuntimeError):
    """Neo4j cleanup failed; the seed was not removed from SQLite or the in-memory store."""


class SeedMaterial(BaseModel):
    """Represents a processed seed material document."""

    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    filename: str
    content_type: str  # "text/plain", "text/markdown", "application/pdf"
    raw_content: str
    processed_content: str | None = None
    status: str = "uploaded"  # uploaded, processing, processed, failed
    entity_count: int = 0
    relationship_count: int = 0
    error_message: str | None = None


class SeedProcessor:
    """Processes seed materials into structured content."""

    def __init__(self):
        self._store = _seed_store
        self._repo: object | None = None

    def _get_repo(self):
        """Lazy-import SeedRepository to avoid circular imports."""
        if self._repo is None:
            from app.database import SeedRepository

            self._repo = SeedRepository()
        return self._repo

    def get_store(self) -> dict[str, SeedMaterial]:
        """Get the seed material store."""
        return self._store

    async def process_file(self, filename: str, content: bytes, content_type: str) -> SeedMaterial:
        """Process an uploaded file into a SeedMaterial."""
        seed = SeedMaterial(
            filename=filename,
            content_type=content_type,
            raw_content="",
        )

        try:
            seed.status = "processing"

            # Extract text based on content type
            if content_type == "application/pdf":
                raw_text = await self.extract_text_from_pdf(content)
            elif content_type in ("text/plain", "text/markdown", "text/text"):
                raw_text = content.decode("utf-8", errors="ignore")
            elif content_type == "application/vnd.ms-powerpoint" or (filename and filename.lower().endswith(".ppt")):
                raise ValueError("Legacy PowerPoint (.ppt) is not supported. Save as .pptx and upload again.")
            elif (content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation") or (
                filename and filename.lower().endswith(".pptx")
            ):
                raw_text = self._extract_text_from_pptx(content)
            elif content_type == "application/vnd.ms-excel" or (filename and filename.lower().endswith(".xls")):
                raise ValueError("Legacy Excel (.xls) is not supported. Save as .xlsx and upload again.")
            elif (content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet") or (
                filename and filename.lower().endswith(".xlsx")
            ):
                raw_text = self._extract_text_from_xlsx(content)
            elif content_type.startswith("image/") or (
                filename and filename.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"))
            ):
                raw_text = self._extract_text_from_image(content, filename or "image")
            elif content_type in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",) or (
                filename and filename.lower().endswith(".docx")
            ):
                raw_text = self._extract_text_from_docx(content)
            else:
                # Try to decode as text for unknown types
                try:
                    raw_text = content.decode("utf-8", errors="ignore")
                except Exception:
                    raise ValueError(f"Unsupported content type: {content_type}")

            seed.raw_content = raw_text
            # Can add additional processing later
            seed.processed_content = raw_text
            seed.status = "processed"

            logger.info(f"Processed {filename} ({content_type}) -> " f"{len(raw_text)} chars")

        except Exception as e:
            seed.status = "failed"
            seed.error_message = str(e)
            logger.error(f"Failed to process file {filename}: {e}")

        # Store in memory and persist to DB
        self._store[seed.id] = seed
        await self._get_repo().save(seed)
        return seed

    @staticmethod
    def _extract_text_from_pptx(content: bytes) -> str:
        """Extract text from PPTX bytes."""
        import io

        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            parts: list[str] = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                slide_texts.append(row_text)
                if slide_texts:
                    parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
            return "\n\n".join(parts) if parts else "[PPTX: no text content found]"
        except ImportError:
            logger.warning("python-pptx not installed; install it for PPTX support")
            return "[PPTX extraction requires python-pptx]"
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            return f"[PPTX extraction failed: {e}]"

    @staticmethod
    def _extract_text_from_xlsx(content: bytes) -> str:
        """Extract text from XLSX bytes."""
        import io

        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            parts: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    line = " | ".join(cells)
                    if line.strip(" |"):
                        rows.append(line)
                if rows:
                    parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            wb.close()
            return "\n\n".join(parts) if parts else "[XLSX: no data found]"
        except ImportError:
            logger.warning("openpyxl not installed; install it for XLSX support")
            return "[XLSX extraction requires openpyxl]"
        except Exception as e:
            logger.error(f"XLSX extraction error: {e}")
            return f"[XLSX extraction failed: {e}]"

    @staticmethod
    def _normalize_exif_text(value: object) -> str | None:
        """Normalize EXIF string/bytes (e.g. UserComment with encoding prefix) to a safe str."""
        if value is None:
            return None
        if isinstance(value, str):
            text = value.strip()
            return text if text else None
        if isinstance(value, bytes):
            raw = value
            prefixes = (
                b"ASCII\x00\x00\x00",
                b"UNICODE\x00",
                b"UNDEFINED\x00",
                b"JIS\x00\x00\x00",
            )
            for p in prefixes:
                if raw.startswith(p):
                    raw = raw[len(p) :]
                    break
            try:
                text = raw.decode("utf-8").strip()
            except UnicodeDecodeError:
                text = raw.decode("latin-1", errors="replace").strip()
            return text if text else None
        try:
            text = str(value).strip()
            return text if text else None
        except Exception:
            return None

    @staticmethod
    def _extract_text_from_image(content: bytes, filename: str) -> str:
        """Extract metadata and description placeholder from image bytes."""
        import io

        try:
            from PIL import Image

            img = Image.open(io.BytesIO(content))
            width, height = img.size
            mode = img.mode
            fmt = img.format or "unknown"
            info_lines = [
                f"[Image: {filename}]",
                f"Format: {fmt}, Size: {width}x{height}, Mode: {mode}",
            ]
            # Extract EXIF text data if present
            exif = img.getexif()
            if exif:
                description = SeedProcessor._normalize_exif_text(exif.get(270, ""))
                if description:
                    info_lines.append(f"Description: {description}")
                user_comment = SeedProcessor._normalize_exif_text(exif.get(37510, ""))
                if user_comment:
                    info_lines.append(f"Comment: {user_comment}")
            img.close()
            return "\n".join(info_lines)
        except ImportError:
            logger.warning("Pillow not installed; install it for image support")
            return f"[Image: {filename}, {len(content)} bytes. Install Pillow for metadata extraction.]"
        except Exception as e:
            logger.error(f"Image extraction error: {e}")
            return f"[Image metadata extraction failed: {e}]"

    @staticmethod
    def _extract_text_from_docx(content: bytes) -> str:
        """Extract text from DOCX bytes using zipfile (no extra deps)."""
        import io
        import re
        import zipfile

        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                if "word/document.xml" not in zf.namelist():
                    return "[DOCX: no document.xml found]"
                xml_content = zf.read("word/document.xml").decode("utf-8", errors="ignore")
                # Strip XML tags, keep text
                text = re.sub(r"<[^>]+>", " ", xml_content)
                lines = [line.strip() for line in text.split("\n") if line.strip()]
                cleaned = "\n".join(lines)
                return cleaned if cleaned else "[DOCX: no text content found]"
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            return f"[DOCX extraction failed: {e}]"

    async def extract_text_from_pdf(self, content: bytes) -> str:
        """Extract text from PDF bytes.

        Uses a simple best-effort approach. For production use,
        consider adding pdfplumber or PyPDF2 dependency.
        """
        try:
            # Try to extract text using basic PDF text extraction
            # PDF files have text embedded between streams
            text_parts = []

            # Simple heuristic: look for text between BT (Begin Text)
            # and ET (End Text)
            # and between stream/endstream
            import re

            # Try to decode as utf-8 first, ignoring errors
            decoded = content.decode("utf-8", errors="ignore")

            # Look for text in parentheses (common PDF text format)
            text_in_parens = re.findall(r"\(([^)]+)\)", decoded)

            # Filter out short strings and binary-looking content
            for text in text_in_parens:
                if len(text) > 3 and all(ord(c) < 128 for c in text):
                    text_parts.append(text)

            if text_parts:
                return " ".join(text_parts)

            # Fallback: try to extract readable ASCII text
            readable_chars = []
            for byte in content:
                # printable ASCII or newline
                if 32 <= byte <= 126 or byte in (10, 13):
                    readable_chars.append(chr(byte))

            extracted = "".join(readable_chars)

            # Clean up: remove excessive whitespace
            lines = [line.strip() for line in extracted.split("\n") if line.strip()]
            cleaned_text = "\n".join(lines)

            if len(cleaned_text) > 100:
                return cleaned_text

            # If extraction yields very little, return a placeholder
            logger.warning("PDF text extraction yielded limited results. " "Consider adding pdfplumber.")
            return f"[PDF: {len(content)} bytes. Basic extraction applied.]\n\n" f"{cleaned_text[:2000]}"

        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            return f"[PDF extraction failed: {str(e)}]"

    async def chunk_content(self, content: str, chunk_size: int = 2000, overlap: int = 200) -> list[str]:
        """Split content into overlapping chunks for processing."""
        if not content:
            return []

        chunks = []
        start = 0
        content_length = len(content)

        while start < content_length:
            end = min(start + chunk_size, content_length)

            # Try to break at a sentence or word boundary
            if end < content_length:
                # Look for sentence endings
                for i in range(min(end + 100, content_length) - 1, start, -1):
                    if content[i] in ".!?":
                        end = i + 1
                        break
                else:
                    # Look for word boundaries
                    for i in range(end, start, -1):
                        if content[i - 1].isspace():
                            end = i
                            break

            chunk = content[start:end].strip()
            if chunk:
                chunks.append(chunk)

            # Move start forward with overlap
            start = end - overlap if end < content_length else content_length

        logger.info(f"Chunked into {len(chunks)} chunks " f"(size={chunk_size}, overlap={overlap})")
        return chunks

    async def get_seed(self, seed_id: str) -> SeedMaterial | None:
        """Get a seed material by ID (in-memory first, then DB)."""
        cached = self._store.get(seed_id)
        if cached is not None:
            return cached
        seed = await self._get_repo().get(seed_id)
        if seed is not None:
            self._store[seed_id] = seed
        return seed

    async def list_seeds(self) -> list[SeedMaterial]:
        """List all seed materials from DB."""
        db_summaries = await self._get_repo().list_all()
        seeds: list[SeedMaterial] = []
        for summary in db_summaries:
            seed_id = summary["id"]
            cached = self._store.get(seed_id)
            if cached is not None:
                seeds.append(cached)
            else:
                seed = await self._get_repo().get(seed_id)
                if seed is not None:
                    self._store[seed_id] = seed
                    seeds.append(seed)
        return seeds

    async def augment_seed(self, seed_id: str) -> SeedMaterial | None:
        """Augment a seed material with external research context.

        Researches entities found in the seed's raw content and appends
        the augmented context to processed_content.

        Args:
            seed_id: ID of the seed material to augment.

        Returns:
            Updated SeedMaterial with augmented content, or None if not found.
        """
        seed = await self.get_seed(seed_id)
        if seed is None:
            logger.warning(f"Seed {seed_id} not found for augmentation")
            return None

        try:
            from app.research.service import research_service

            result = await research_service.augment_text(seed.raw_content)

            augmented_context = result.get("augmented_context", "")
            if augmented_context:
                separator = "\n\n--- AUTORESEARCH CONTEXT ---\n\n"
                base_content = seed.processed_content or seed.raw_content
                seed.processed_content = base_content + separator + augmented_context

            seed.status = "augmented"

            self._store[seed.id] = seed
            await self._get_repo().save(seed)

            entities_found = result.get("entities_found", [])
            logger.info(f"Augmented seed {seed_id} with research context " f"({len(entities_found)} entities found)")

        except Exception as e:
            logger.error(f"Failed to augment seed {seed_id}: {e}")
            seed.status = "failed"
            seed.error_message = f"Augmentation failed: {e}"
            self._store[seed.id] = seed
            await self._get_repo().save(seed)

        return seed

    async def delete_seed(self, seed_id: str) -> bool:
        """Delete a seed material from memory, DB, and any Neo4j graph scoped to this seed."""
        in_memory = seed_id in self._store
        if not in_memory and (await self._get_repo().get(seed_id)) is None:
            return False

        tombstone_seed_graph(seed_id)
        neo4j = get_application_neo4j_client()
        if neo4j is not None:
            try:
                await neo4j.clear_graph(seed_id)
            except Exception as e:
                logger.exception(
                    "Neo4j cleanup failed for seed %s; seed not deleted from storage",
                    seed_id,
                )
                clear_seed_graph_tombstone(seed_id)
                raise SeedGraphCleanupError(f"Neo4j graph cleanup failed for seed {seed_id}") from e
        self._store.pop(seed_id, None)
        return await self._get_repo().delete(seed_id)

    async def update_seed(self, seed_id: str, **updates) -> SeedMaterial | None:
        """Update a seed material in memory and DB."""
        seed = await self.get_seed(seed_id)
        if seed is None:
            return None

        for key, value in updates.items():
            if hasattr(seed, key):
                setattr(seed, key, value)

        self._store[seed_id] = seed
        await self._get_repo().save(seed)
        return seed
